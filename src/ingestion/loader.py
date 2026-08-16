import os
import io
from typing import List, Set
import fitz  # PyMuPDF
from langchain_core.documents import Document
from src.chunking.chunker import split_documents
from src.embeddings.embedder import get_embedding_model
from src.vectordb.vector_store import get_vector_store, VectorStoreManager
from src.utils.helpers import load_app_config


def is_image_pdf(pdf_doc: fitz.Document) -> bool:
    """Check if PDF is primarily scanned images without selectable text."""
    total_text_len = 0
    total_images = 0
    for page in pdf_doc:
        total_text_len += len(page.get_text().strip())
        total_images += len(page.get_images())
    return total_text_len < 100 and total_images > 0


def load_single_pdf(file_path: str) -> List[Document]:
    """Extract text from PDF file with automatic OCR fallback for scanned pages."""
    docs = []
    pdf = fitz.open(file_path)

    if is_image_pdf(pdf):
        try:
            import pytesseract
            from PIL import Image

            # Default Windows Tesseract path if available
            tesseract_default = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
            if os.path.exists(tesseract_default):
                pytesseract.pytesseract.tesseract_cmd = tesseract_default

            print(f"  [OCR SCANNING] {os.path.basename(file_path)} ({len(pdf)} pages)")
            for page_num, page in enumerate(pdf):
                print(f"    → Page {page_num+1}/{len(pdf)}", end="\r")
                pix = page.get_pixmap(dpi=150)
                img = Image.open(io.BytesIO(pix.tobytes("png")))
                text = pytesseract.image_to_string(img)
                if text.strip():
                    docs.append(Document(
                        page_content=text,
                        metadata={"source": file_path, "page": page_num + 1, "type": "ocr_scanned"}
                    ))
            print()
        except Exception as e:
            print(f"\nOCR Extraction Warning for {file_path}: {e}")
    else:
        for page_num, page in enumerate(pdf):
            text = page.get_text()
            if text.strip():
                docs.append(Document(
                    page_content=text,
                    metadata={"source": file_path, "page": page_num + 1, "type": "digital_pdf"}
                ))
    pdf.close()
    return docs


def load_single_docx(file_path: str) -> List[Document]:
    """Extract text from DOCX file."""
    try:
        import docx
        doc = docx.Document(file_path)
        full_text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        if full_text.strip():
            return [Document(page_content=full_text, metadata={"source": file_path, "type": "docx"})]
    except Exception as e:
        print(f"  ⚠️ Skipping unreadable DOCX {os.path.basename(file_path)}: {e}")
    return []


def load_single_pptx(file_path: str) -> List[Document]:
    """Extract text from PPTX slides."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        docs = []
        for idx, slide in enumerate(prs.slides):
            text = "\n".join([
                shape.text for shape in slide.shapes
                if shape.has_text_frame and shape.text.strip()
            ])
            if text.strip():
                docs.append(Document(
                    page_content=text,
                    metadata={"source": file_path, "slide": idx + 1, "type": "pptx"}
                ))
        return docs
    except Exception as e:
        print(f"  ⚠️ Skipping unreadable PPTX {os.path.basename(file_path)}: {e}")
    return []


def load_document(file_path: str) -> List[Document]:
    """Universal document loader dispatching by extension."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return load_single_pdf(file_path)
    elif ext == ".docx":
        return load_single_docx(file_path)
    elif ext == ".pptx":
        return load_single_pptx(file_path)
    return []


def process_directory_incrementally(root_path: str = "./PDF_Data"):
    """
    Incrementally ingest documents from folder into the ChromaDB vector database.
    Skips files that are already indexed in ChromaDB.
    """
    cfg = load_app_config()
    db_manager = get_vector_store()
    embedder = get_embedding_model()

    indexed_sources = db_manager.get_indexed_sources()
    print(f"Already indexed files in database: {len(indexed_sources)}")

    supported_extensions = {".pdf", ".docx", ".pptx"}
    all_files = []
    for root, dirs, files in os.walk(root_path):
        for file in files:
            ext = os.path.splitext(file)[1].lower()
            if ext in supported_extensions:
                all_files.append(os.path.join(root, file))

    total_files = len(all_files)
    print(f"Found {total_files} total supported documents in '{root_path}'.\n")

    for idx, file_path in enumerate(all_files, start=1):
        normalized_path = file_path.replace("\\", "/")
        if file_path in indexed_sources or normalized_path in indexed_sources:
            print(f"[{idx}/{total_files}] ⏩ SKIP (Already Indexed): {os.path.basename(file_path)}")
            continue

        print(f"[{idx}/{total_files}] 🔄 Processing: {os.path.basename(file_path)}")
        try:
            docs = load_document(file_path)
            if docs:
                chunks = split_documents(docs, root_data_dir=root_path)
                if chunks:
                    texts = [c.page_content for c in chunks]
                    embeddings = embedder.encode(texts)
                    db_manager.add_documents(chunks, embeddings)
                    print(f"  ✅ Saved {len(chunks)} chunks to vector database.")
            else:
                print(f"  ⚠️ No text extracted from {os.path.basename(file_path)}.")
        except Exception as e:
            print(f"  ❌ Error processing file {file_path}: {e}")

    print("\n🎉 Pipeline Execution Completed! Total chunks in DB:", db_manager.collection.count())
