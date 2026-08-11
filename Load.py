import os, uuid
from typing import List, Set
import chromadb
import numpy as np
import fitz  # PyMuPDF
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from sentence_transformers import SentenceTransformer


# ─────────────────────────────────────────────────────────────
# EMBEDDING MODEL
# ─────────────────────────────────────────────────────────────

class EmbeddingM:
    def __init__(self, model_N: str = "all-MiniLM-L6-v2"):
        self.model = None
        self.model_N = model_N
        self._load_model()

    def _load_model(self):
        try:
            self.model = SentenceTransformer(self.model_N)
            print(f"Model Loaded successfully. Embedding Model: {self.model_N}")
        except Exception as e:
            print(f"Error loading model {self.model_N}: {e}")
            raise

    def split_texts(self, docs: List[Document], chunk_S: int = 1000, chunk_overlap: int = 200) -> List[Document]:
        text_S = RecursiveCharacterTextSplitter(chunk_size=chunk_S, chunk_overlap=chunk_overlap)
        split_docs: List[Document] = []
        for doc in docs:
            chunks = text_S.split_text(doc.page_content)
            for chunk in chunks:
                source_path = doc.metadata.get("source", "")
                parts = source_path.replace("\\", "/").split("/")
                subject_N = parts[1] if len(parts) > 2 else "General"
                metadata = {
                    **doc.metadata,
                    "source": source_path,
                    "subject": subject_N
                }
                split_docs.append(Document(page_content=chunk, metadata=metadata))
        return split_docs

    def generate_E(self, texts: List[str]) -> np.ndarray:
        if not self.model:
            raise ValueError("Model is not loaded.")
        if isinstance(texts, str):
            texts = [texts]
        e = self.model.encode(texts, show_progress_bar=False)
        return e


# ─────────────────────────────────────────────────────────────
# VECTOR STORE (ChromaDB with Checkpoint Support)
# ─────────────────────────────────────────────────────────────

class VectoreS:
    def __init__(self, collection_N: str = "Document__C", persist_D: str = "./pdf_db/chromadb"):
        self.collection_N = collection_N
        self.persistent_D = persist_D
        self.client = None
        self.collection = None
        self._initialize_S()

    def _initialize_S(self):
        try:
            os.makedirs(self.persistent_D, exist_ok=True)
            self.client = chromadb.PersistentClient(path=self.persistent_D)
            self.collection = self.client.get_or_create_collection(
                name=self.collection_N,
                metadata={"description": "PDF embedding document for RAG"}
            )
            print(f"Vector Store initialized : Collection '{self.collection_N}'")
            print(f"Current document count in collection : {self.collection.count()}")
        except Exception as e:
            print(f"Error while initializing Vector Store : {e}")
            raise

    def get_indexed_sources(self) -> Set[str]:
        """Fetch all unique source file paths already stored in ChromaDB."""
        try:
            results = self.collection.get(include=["metadatas"])
            metadatas = results.get("metadatas", [])
            indexed_sources = {m.get("source") for m in metadatas if m and "source" in m}
            return indexed_sources
        except Exception as e:
            print(f"Warning: Could not fetch indexed sources: {e}")
            return set()

    def add_Document(self, documents: List[Document], embeddings: np.ndarray, batch_size: int = 500):
        if len(documents) != len(embeddings):
            raise ValueError("Number of documents must match number of embeddings")
        if len(documents) == 0:
            return

        for start in range(0, len(documents), batch_size):
            end = min(start + batch_size, len(documents))
            batch_docs = documents[start:end]
            batch_embeddings = embeddings[start:end]
            ids, metadatas, doc_T, embedding_L = [], [], [], []
            
            for i, (doc, embedd) in enumerate(zip(batch_docs, batch_embeddings), start=start):
                d_id = f"doc_{uuid.uuid4().hex[:8]}_{i}"
                ids.append(d_id)
                metadata = dict(doc.metadata or {})
                metadata["content_length"] = len(doc.page_content)
                metadatas.append(metadata)
                doc_T.append(doc.page_content)
                embedding_L.append(embedd.tolist())

            try:
                self.collection.add(ids=ids, embeddings=embedding_L, metadatas=metadatas, documents=doc_T)
            except Exception as e:
                print(f"  ❌ Error adding batch to vector store: {e}")


# ─────────────────────────────────────────────────────────────
# SMART MULTI-FORMAT DOCUMENT LOADERS
# ─────────────────────────────────────────────────────────────

def is_image_pdf(pdf_doc: fitz.Document) -> bool:
    """Check if PDF has sparse text and page images (scanned PDF)."""
    total_text_len = 0
    total_images = 0
    for page in pdf_doc:
        total_text_len += len(page.get_text().strip())
        total_images += len(page.get_images())
    
    # If text is minimal (<100 chars total across doc) and contains image objects
    return total_text_len < 100 and total_images > 0


def load_single_pdf(file_path: str) -> List[Document]:
    docs = []
    pdf = fitz.open(file_path)

    if is_image_pdf(pdf):
        try:
            import pytesseract
            from PIL import Image
            import io

            pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
            print(f"  [OCR SCANNING] {os.path.basename(file_path)} ({len(pdf)} pages)")
            
            for page_num, page in enumerate(pdf):
                print(f"    → Page {page_num + 1}/{len(pdf)}", end="\r")
                pix = page.get_pixmap(dpi=150)  # Optimized 150 DPI
                img = Image.open(io.BytesIO(pix.tobytes("png")))
                text = pytesseract.image_to_string(img)
                if text.strip():
                    docs.append(Document(
                        page_content=text,
                        metadata={"source": file_path, "page": page_num + 1, "type": "ocr_scanned"}
                    ))
            print() # Clear line return
        except Exception as e:
            print(f"\n  ❌ Tesseract OCR Failed for {file_path}: {e}")
    else:
        for page_num, page in enumerate(pdf):
            text = page.get_text()
            if text.strip():
                docs.append(Document(
                    page_content=text,
                    metadata={"source": file_path, "page": page_num + 1, "type": "digital_pdf"}
                ))
    pdf.close()
    return docs


def load_single_docx(file_path: str) -> List[Document]:
    try:
        import docx
        doc = docx.Document(file_path)
        full_text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        if full_text.strip():
            return [Document(page_content=full_text, metadata={"source": file_path, "type": "docx"})]
    except Exception as e:
        print(f"  ⚠️ Skipping corrupt or unreadable DOCX {os.path.basename(file_path)}: {e}")
    return []


def load_single_pptx(file_path: str) -> List[Document]:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        docs = []
        for idx, slide in enumerate(prs.slides):
            text = "\n".join([
                shape.text for shape in slide.shapes
                if shape.has_text_frame and shape.text.strip()
            ])
            if text.strip():
                docs.append(Document(
                    page_content=text,
                    metadata={"source": file_path, "slide": idx + 1, "type": "pptx"}
                ))
        return docs
    except Exception as e:
        print(f"  ⚠️ Skipping corrupt PPTX {os.path.basename(file_path)}: {e}")
    return []


# ─────────────────────────────────────────────────────────────
# INCREMENTAL PIPELINE EXECUTION
# ─────────────────────────────────────────────────────────────

def process_documents_incrementally(root_path: str = "./PDF_Data"):
    embedding_M = EmbeddingM()
    vector_store = VectoreS()
    
    # 1. Fetch set of already indexed file paths to enable resuming
    indexed_sources = vector_store.get_indexed_sources()
    print(f"Already indexed files in database: {len(indexed_sources)}")

    # 2. Gather all candidate files
    supported_extensions = {".pdf", ".docx", ".pptx"}
    all_files = []
    for root, dirs, files in os.walk(root_path):
        for file in files:
            ext = os.path.splitext(file)[1].lower()
            if ext in supported_extensions:
                all_files.append(os.path.join(root, file))

    total_files = len(all_files)
    print(f"Found {total_files} total supported documents in '{root_path}'.\n")

    # 3. Process file by file
    for idx, file_path in enumerate(all_files, start=1):
        normalized_path = file_path.replace("\\", "/")
        
        # Checkpoint check
        if file_path in indexed_sources or normalized_path in indexed_sources:
            print(f"[{idx}/{total_files}] ⏩ SKIP (Already Indexed): {os.path.basename(file_path)}")
            continue

        print(f"[{idx}/{total_files}] 🔄 Processing: {os.path.basename(file_path)}")
        ext = os.path.splitext(file_path)[1].lower()
        
        try:
            if ext == ".pdf":
                docs = load_single_pdf(file_path)
            elif ext == ".docx":
                docs = load_single_docx(file_path)
            elif ext == ".pptx":
                docs = load_single_pptx(file_path)
            else:
                docs = []

            if docs:
                chunks = embedding_M.split_texts(docs)
                embeddings = embedding_M.generate_E([c.page_content for c in chunks])
                vector_store.add_Document(chunks, embeddings)
                print(f"  ✅ Saved {len(chunks)} chunks to vector database.")
            else:
                print(f"  ⚠️ No text extracted from {os.path.basename(file_path)}.")
                
        except Exception as e:
            print(f"  ❌ Error processing file {file_path}: {e}")

    print("\n🎉 Pipeline Execution Completed! Total chunks in DB:", vector_store.collection.count())


if __name__ == "__main__":
    process_documents_incrementally("./PDF_Data")
